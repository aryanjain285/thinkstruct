"""Generate the presentation deck as a .pptx.

    python scripts/make_deck.py [-o docs/thinkstruct-deck.pptx]

Kept as a script rather than a binary blob so the deck regenerates when the numbers do.
"""
from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------- identity
INK = RGBColor(0x13, 0x1A, 0x24)
INK2 = RGBColor(0x35, 0x42, 0x4F)
MUTED = RGBColor(0x61, 0x70, 0x7F)
DRAFT = RGBColor(0x1C, 0x4F, 0xD8)
DRAFT_SOFT = RGBColor(0xE6, 0xEC, 0xFD)
SIGNAL = RGBColor(0xA4, 0x53, 0x0B)
SIGNAL_SOFT = RGBColor(0xFB, 0xF0, 0xE2)
GOOD = RGBColor(0x15, 0x68, 0x3C)
GOOD_SOFT = RGBColor(0xE3, 0xF2, 0xE9)
PAPER = RGBColor(0xF4, 0xF6, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RULE = RGBColor(0xD5, 0xDB, 0xE2)

SERIF = "Georgia"
SANS = "Calibri"
MONO = "Consolas"

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.62)                      # page margin
CONTENT_W = W - 2 * M


def textbox(slide, x, y, w, h, *, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, *, size=14, bold=False, color=INK2, font=SANS, space_after=6,
         space_before=0, align=PP_ALIGN.LEFT, first=False, line=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    if line:
        p.line_spacing = line
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return p


def rich(tf, parts, *, size=14, space_after=6, space_before=0, first=False, line=1.25):
    """parts: list of (text, bold, color, font)."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.line_spacing = line
    for text, bold, color, font in parts:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return p


def box(slide, x, y, w, h, *, fill=WHITE, line=RULE, width=Pt(0.75),
        shape=MSO_SHAPE.RECTANGLE, shadow=False):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = width
    s.shadow.inherit = shadow
    s.text_frame.word_wrap = True
    return s


def label_box(slide, x, y, w, h, title, subtitle=None, *, fill=WHITE, line=RULE,
              title_size=12, sub_size=9, title_color=INK, accent=None):
    """A diagram node: title, optional subtitle, optional accent stripe."""
    s = box(slide, x, y, w, h, fill=fill, line=line)
    if accent:
        stripe = box(slide, x, y, Emu(int(Inches(0.055))), h, fill=accent, line=None)
        stripe.shadow.inherit = False
    tf = s.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = 0
    para(tf, title, size=title_size, bold=True, color=title_color, font=SANS,
         space_after=0, align=PP_ALIGN.CENTER, first=True)
    if subtitle:
        para(tf, subtitle, size=sub_size, color=MUTED, font=MONO,
             space_after=0, align=PP_ALIGN.CENTER)
    return s


def arrow(slide, x, y, w, *, color=DRAFT, thickness=Pt(1.25)):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, Inches(0.14))
    a.fill.solid()
    a.fill.fore_color.rgb = color
    a.line.fill.background()
    a.shadow.inherit = False
    return a


def down_arrow(slide, x, y, h, *, color=DRAFT):
    a = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x, y, Inches(0.14), h)
    a.fill.solid()
    a.fill.fore_color.rgb = color
    a.line.fill.background()
    a.shadow.inherit = False
    return a


def hrule(slide, x, y, w, color=RULE):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(0.75))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def new_slide(prs, n=None, eyebrow=None, title=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = box(s, 0, 0, W, H, fill=WHITE, line=None)
    bg.shadow.inherit = False
    y = M
    if n is not None:
        tab = box(s, 0, 0, Inches(0.62), Inches(0.32), fill=DRAFT, line=None)
        tab.shadow.inherit = False
        tf = tab.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        para(tf, n, size=10, bold=True, color=WHITE, font=MONO,
             space_after=0, align=PP_ALIGN.CENTER, first=True)
    if eyebrow:
        tf = textbox(s, M, y, CONTENT_W, Inches(0.24))
        para(tf, eyebrow.upper(), size=10, bold=True, color=DRAFT, font=MONO,
             space_after=0, first=True)
        y += Inches(0.3)
    if title:
        # Georgia bold at 30pt fits roughly 54 characters across the content width;
        # anything longer wraps and needs a second line reserved.
        two_line = len(title) > 52
        h = Inches(1.22) if two_line else Inches(0.62)
        tf = textbox(s, M, y, CONTENT_W, h)
        para(tf, title, size=30, bold=True, color=INK, font=SERIF,
             space_after=0, first=True, line=1.02)
        y += Inches(1.32) if two_line else Inches(0.78)
    return s, y


def stat(slide, x, y, w, number, label):
    bar = box(slide, x, y, Pt(2.5), Inches(0.86), fill=DRAFT, line=None)
    bar.shadow.inherit = False
    tf = textbox(slide, x + Inches(0.12), y, w - Inches(0.12), Inches(0.86))
    para(tf, number, size=22, bold=True, color=INK, font=MONO, space_after=2, first=True)
    para(tf, label, size=10, color=MUTED, font=SANS, space_after=0, line=1.15)


def callout(slide, x, y, w, h, heading, body, *, kind="signal"):
    fill, edge = {
        "signal": (SIGNAL_SOFT, SIGNAL),
        "good": (GOOD_SOFT, GOOD),
        "draft": (DRAFT_SOFT, DRAFT),
    }[kind]
    box(slide, x, y, w, h, fill=fill, line=None)
    stripe = box(slide, x, y, Pt(3), h, fill=edge, line=None)
    stripe.shadow.inherit = False
    tf = textbox(slide, x + Inches(0.16), y + Inches(0.12), w - Inches(0.3), h - Inches(0.2))
    if heading:
        para(tf, heading, size=12, bold=True, color=INK, font=SANS, space_after=4, first=True)
        para(tf, body, size=11, color=INK2, font=SANS, space_after=0, line=1.2)
    else:
        para(tf, body, size=11, color=INK2, font=SANS, space_after=0, first=True, line=1.2)


def table(slide, x, y, w, headers, rows, *, col_w=None, highlight=None,
          font_size=11, row_h=Inches(0.3)):
    nr, nc = len(rows) + 1, len(headers)
    shape = slide.shapes.add_table(nr, nc, x, y, w, row_h * nr)
    tbl = shape.table
    tbl.first_row = False
    if col_w:
        total = sum(col_w)
        for i, frac in enumerate(col_w):
            tbl.columns[i].width = Emu(int(w * frac / total))
    for j, htxt in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = ""
        c.fill.solid()
        c.fill.fore_color.rgb = WHITE
        c.margin_left = c.margin_right = Inches(0.07)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        para(c.text_frame, htxt.upper(), size=8.5, bold=True, color=MUTED, font=MONO,
             space_after=0, first=True,
             align=PP_ALIGN.RIGHT if j and htxt.startswith("~") is False and j > 0 and headers[j][:1].isupper() is False else PP_ALIGN.LEFT)
    for i, row in enumerate(rows, start=1):
        hot = highlight is not None and i - 1 == highlight
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.text = ""
            c.fill.solid()
            c.fill.fore_color.rgb = DRAFT_SOFT if hot else WHITE
            c.margin_left = c.margin_right = Inches(0.07)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            is_num = j > 0 and any(ch.isdigit() for ch in val)
            para(c.text_frame, val,
                 size=font_size, bold=hot, color=INK if hot else INK2,
                 font=MONO if is_num else SANS, space_after=0, first=True,
                 align=PP_ALIGN.RIGHT if is_num else PP_ALIGN.LEFT)
    return tbl


def mono_block(slide, x, y, w, h, lines):
    box(slide, x, y, w, h, fill=PAPER, line=RULE)
    tf = textbox(slide, x + Inches(0.16), y + Inches(0.12), w - Inches(0.3), h - Inches(0.2))
    for i, (text, color) in enumerate(lines):
        para(tf, text, size=10, color=color, font=MONO, space_after=1,
             first=(i == 0), line=1.15)


# ================================================================== slides
def notes(slide, text):
    """Q&A prep lives in the speaker notes — visible in presenter view, not to the room."""
    slide.notes_slide.notes_text_frame.text = text.strip()


def slide_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = box(s, 0, 0, W, H, fill=WHITE, line=None); bg.shadow.inherit = False
    band = box(s, 0, 0, Inches(0.15), H, fill=DRAFT, line=None); band.shadow.inherit = False
    tab = box(s, Inches(0.15), 0, Inches(0.6), Inches(0.32), fill=DRAFT, line=None)
    tab.shadow.inherit = False
    tftab = tab.text_frame; tftab.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tftab, "01", size=10, bold=True, color=WHITE, font=MONO, space_after=0,
         align=PP_ALIGN.CENTER, first=True)

    L = Inches(0.85)
    CW = W - L - M
    y = Inches(0.62)

    tf = textbox(s, L, y, CW, Inches(0.26))
    para(tf, "THINKSTRUCT · CODING TASK · ARYAN JAIN", size=10, bold=True, color=DRAFT,
         font=MONO, space_after=0, first=True)
    y += Inches(0.34)

    tf = textbox(s, L, y, CW, Inches(1.45))
    para(tf, "Finding the claim that already covers your invention",
         size=32, bold=True, color=INK, font=SERIF, space_after=0, first=True, line=1.06)
    y += Inches(1.48)

    tf = textbox(s, L, y, Inches(11.0), Inches(0.6))
    rich(tf, [("Before filing, litigating or committing R&D budget someone must answer: ",
               False, INK2, SANS),
              ("has anyone already claimed this?", True, INK, SANS),
              ("  A ", False, INK2, SANS),
              ("carbon fibre spoke", True, INK, SANS),
              (" may be claimed elsewhere as an ", False, INK2, SANS),
              ("elongate composite tension member", True, INK, SANS),
              (" — keyword search misses it, and reading 10M patents is not an option.",
               False, INK2, SANS)],
         size=13.5, first=True, space_after=0, line=1.35)
    y += Inches(0.82)

    cw = (CW - Inches(0.3)) / 2
    callout(s, L, y, cw, Inches(1.05), "Missing prior art",
            "A patent issues that is invalid. It surfaces in litigation, where "
            "invalidation runs $1–3M and the R&D is already sunk.")
    callout(s, L + cw + Inches(0.3), y, cw, Inches(1.05), "Returning noise",
            "Attorneys read hundreds of irrelevant patents at $400–700/hour. 200 vague "
            "hits has saved nobody anything.")
    y += Inches(1.22)

    hrule(s, L, y, CW); y += Inches(0.22)
    tf = textbox(s, L, y, CW, Inches(0.26))
    para(tf, "DESIGN DECISION · RETRIEVAL OPERATES ON CLAIMS, NEVER WHOLE PATENTS",
         size=10, bold=True, color=DRAFT, font=MONO, space_after=0, first=True)
    y += Inches(0.4)

    sw = (CW - Inches(0.6)) / 3
    stat(s, L, y, sw, "16.5", "average claims per patent — one blob matches every query "
                              "in the field")
    stat(s, L + sw + Inches(0.3), y, sw, "Legal",
         "infringement is decided claim by claim; “patent X is relevant” is not actionable")
    stat(s, L + 2 * (sw + Inches(0.3)), y, sw, "Evidence",
         "every result names the matching record, so output is checked not trusted")
    y += Inches(1.0)

    callout(s, L, y, CW, Inches(0.66), None,
            "Scope boundary — the system surfaces and ranks technically similar claims. It "
            "does not output novelty, validity or infringement conclusions; those are legal "
            "determinations, and the reranker prompt says so explicitly.", kind="draft")

    tf = textbox(s, L, Inches(7.02), CW, Inches(0.26))
    para(tf, "640 patents · 18,743 indexed records · 329 tests · "
             "github.com/aryanjain285/thinkstruct",
         size=10, color=MUTED, font=MONO, space_after=0, first=True)

    notes(s, """
HOW DO YOU HANDLE MISSING FIELDS?
Measured: no field is ever absent, and none is ever empty — except descriptions. 119
patents (18.6%) have descriptions that are present but entirely blank paragraphs. Those
are KEPT and FLAGGED, never excluded: validation grades them 'warning', not 'error',
because they stay fully searchable via claims and abstract. Excluding them would discard
18.6% of the corpus for no gain. Nothing in this dataset is dropped.

WHY CLAIM-LEVEL AND NOT PATENT-LEVEL?
Three reasons. (1) Claims are the legal unit — infringement and novelty are determined
claim by claim, so "patent X is relevant" is not an actionable answer. (2) Precision — a
46-claim patent indexed as one document matches nearly any query in its field. (3)
Evidence — every result names the record that matched, so it can be verified.
Results are regrouped by patent for display, scored max(record) + 0.3*mean(next 3), so a
patent matching on several claims outranks one matching by fluke on a single passage.

WHY NOT JUST USE AN LLM OVER THE WHOLE CORPUS?
15.4M characters is ~3.8M tokens per query. At scale that is 10M patents. Retrieval
narrows to 50 candidates first; the LLM only ever sees those.
""")


def slide_architecture(prs):
    s, y = new_slide(prs, "02", "Architecture",
                     "Two retrievers, one filter set, fused and reranked")

    # ---- pipeline row
    nodes = [("Raw JSON", "64 files"), ("Validate", "graded"), ("Normalise", "NFKC"),
             ("Reconstruct", "claims"), ("Records", "18,743")]
    nw, gap = Inches(1.72), Inches(0.42)
    x = M
    for i, (t, sub) in enumerate(nodes):
        accent = SIGNAL if t == "Reconstruct" else None
        label_box(s, x, y, nw, Inches(0.62), t, sub, accent=accent)
        if i < len(nodes) - 1:
            arrow(s, x + nw + Inches(0.08), y + Inches(0.24), gap - Inches(0.16))
        x += nw + gap
    ix = x
    label_box(s, ix, y, Inches(1.45), Inches(0.62), "OpenSearch", "BM25 + kNN",
              fill=DRAFT_SOFT, line=DRAFT)
    y += Inches(0.92)

    # ---- query row
    down_arrow(s, M + Inches(0.75), y - Inches(0.24), Inches(0.2))
    qy = y + Inches(0.1)
    label_box(s, M, qy, Inches(1.6), Inches(0.55), "Query", None)
    arrow(s, M + Inches(1.68), qy + Inches(0.2), Inches(0.34))
    label_box(s, M + Inches(2.1), qy, Inches(1.85), Inches(0.55), "Pre-filters",
              "CPC · title · abstract", accent=DRAFT)

    bx = M + Inches(4.3)
    label_box(s, bx, qy - Inches(0.34), Inches(1.5), Inches(0.5), "BM25", "lexical")
    label_box(s, bx, qy + Inches(0.34), Inches(1.5), Inches(0.5), "k-NN", "semantic")
    arrow(s, M + Inches(3.98), qy + Inches(0.2), Inches(0.28))

    fx = bx + Inches(1.75)
    arrow(s, fx - Inches(0.22), qy + Inches(0.2), Inches(0.2))
    label_box(s, fx, qy, Inches(1.4), Inches(0.55), "RRF fusion", "scale-free",
              fill=DRAFT_SOFT, line=DRAFT)

    rx = fx + Inches(1.62)
    arrow(s, rx - Inches(0.2), qy + Inches(0.2), Inches(0.18))
    label_box(s, rx, qy, Inches(1.65), Inches(0.55), "Rerank", "LTR / CE / LLM")

    gx = rx + Inches(1.87)
    arrow(s, gx - Inches(0.2), qy + Inches(0.2), Inches(0.18))
    label_box(s, gx, qy, Inches(1.55), Inches(0.55), "Group", "by patent")
    y = qy + Inches(1.0)

    cw = (CONTENT_W - Inches(0.3)) / 2
    callout(s, M, y, cw, Inches(0.95), "Why both retrievers",
            "BM25 nails exact component names, materials and legal phrasing. Vectors "
            "catch the same mechanism in different words.", kind="draft")
    callout(s, M + cw + Inches(0.3), y, cw, Inches(0.95), "Why RRF, not score addition",
            "BM25 scores sit in the tens, cosine in 0–1. RRF is scale-free and needs no "
            "tuned weight.", kind="draft")
    y += Inches(1.15)

    callout(s, M, y, CONTENT_W, Inches(0.68), None,
            "Filters are pre-filters applied inside k-NN graph traversal. Post-filtering "
            "would return far fewer than k under a selective constraint. Verified: with "
            "B60B applied, dense and hybrid still return a full 50/50.", kind="good")
    y += Inches(0.84)

    tf = textbox(s, M, y, CONTENT_W, Inches(0.9))
    para(tf, "EVERY MODEL IS SWAPPABLE BY FLAG — NO CODE CHANGE", size=10, bold=True,
         color=DRAFT, font=MONO, space_after=6, first=True)
    for t in ["Embeddings — all-MiniLM / bge / e5 / gte local, or OpenAI hosted",
              "Reranker — trained LTR, HuggingFace cross-encoder, LLM, or off",
              "Vocabulary — config/synonyms.txt, edit and reindex"]:
        para(tf, "— " + t, size=11.5, color=INK2, font=SANS, space_after=3, line=1.15)

    notes(s, """
WHY OPENSEARCH? WHAT WERE THE ALTERNATIVES?
The requirement is lexical scoring, vector similarity AND metadata filtering resolved in
a single query with shared pre-filters. OpenSearch does all three natively. Splitting
across a keyword store and a vector DB means two round-trips, two filter implementations
that can disagree, and client-side merging.

- Elasticsearch: near-identical capability. OpenSearch is the Apache-2.0 fork, so no
  licence question.
- Postgres + pgvector: excellent if you already run Postgres, but BM25 is bolted on via
  tsvector and weaker on long technical text.
- Pinecone / Weaviate / Qdrant: better pure-vector ergonomics, but you still need a
  lexical engine beside them — and BM25 proved a STRONG baseline here, so dropping it
  would have hurt.
- Vespa: arguably the best technical fit for hybrid ranking, but far heavier to operate
  for a proof of concept.
- FAISS: a library, not a search engine. No filters, no persistence, no lexical side.

The k-NN mapping uses the LUCENE engine specifically because it supports filtering DURING
graph traversal. Every query here carries metadata constraints, so that is not a detail.

WHY RRF RATHER THAN WEIGHTED SCORE COMBINATION?
BM25 scores are unbounded and land in the tens; cosine similarity is bounded 0-1. Adding
them requires normalisation that drifts as the corpus changes. RRF uses only rank
position, so it is scale-free and has no weight to tune. The trade-off: it discards score
magnitude, so a runaway-best match is treated the same as a marginal rank-1.

WHY DENORMALISE THE ABSTRACT ONTO EVERY RECORD?
So title/abstract keyword filters resolve in one pass instead of a patent-level join.
Costs about 20% index size, removes a round-trip. This is the "make hybrid search help
performance" hint in the brief — trade space for filter latency.
""")


def slide_reconstruction(prs):
    s, y = new_slide(prs, "03", "The hard problem",
                     "The claims field is broken, and the obvious fix corrupts it")

    sw = (CONTENT_W - Inches(0.6)) / 3
    stat(s, M, y, sw, "68.9%", "of entries begin with a claim number")
    stat(s, M + sw + Inches(0.3), y, sw, "88.6%", "of patents open with a numberless fragment")
    stat(s, M + 2 * (sw + Inches(0.3)), y, sw, "79.7%", "have no “1.” entry at all")
    y += Inches(0.95)

    tf = textbox(s, M, y, Inches(11.2), Inches(0.4))
    rich(tf, [("Preambles are ", False, INK2, SANS),
              ("stripped, not split", True, INK, SANS),
              (".  Patent 20240051333 (“SPOKE”) opens mid-sentence — the "
               "“1 . A spoke comprising:” is simply gone.", False, INK2, SANS)],
         size=12.5, first=True, space_after=0)
    y += Inches(0.5)

    callout(s, M, y, CONTENT_W, Inches(0.78), "Why the obvious rule fails",
            "“Append every numberless fragment to the previous claim” breaks twice: "
            "nothing precedes an index-0 fragment, so claim 1 is dropped; and mid-list "
            "fragments are often new independent claims, glued onto their predecessor.")
    y += Inches(0.94)

    mono_block(s, M, y, CONTENT_W, Inches(0.86), [
        ("20240051338, entry 1:", MUTED),
        ("  prev: '1 - 5 . (canceled)'", INK2),
        ("  FRAG: 'a tread portion extending in a tire circumferential direction...'", INK2),
        ("                                                    ^^ this is claim 6", SIGNAL),
    ])
    y += Inches(1.02)

    cw = (CONTENT_W - Inches(0.3)) / 2
    tf = textbox(s, M, y, cw, Inches(1.3))
    para(tf, "The rule used instead", size=13, bold=True, color=INK, font=SANS,
         space_after=5, first=True)
    rich(tf, [("Segment first, number second. ", True, INK, SANS),
              ("A fragment starts a new claim if it is first, or if the previous entry "
               "closed on a full stop. Numbers resolve afterwards by anchoring to the "
               "next explicit number. Every uncertain decision is recorded on the claim, "
               "so the UI shows provenance instead of pretending the data was clean.",
               False, INK2, SANS)], size=11.5, space_after=0)

    callout(s, M + cw + Inches(0.3), y, cw, Inches(0.82), None,
            "Invariant, checked over the whole corpus: all 10,578 source entries are "
            "consumed exactly once — zero dropped, zero duplicated.", kind="good")
    tf = textbox(s, M + cw + Inches(0.42), y + Inches(0.95), cw - Inches(0.24), Inches(0.6))
    para(tf, "Tests also cover formats absent from this dataset — “1)” numbering, "
             "em-dash cancellation ranges, 120-claim patents — so the parser is not "
             "fitted to one sample.",
         size=11, color=INK2, font=SANS, space_after=0, first=True, line=1.2)

    notes(s, """
WHY IS THIS THE MOST IMPORTANT SLIDE?
Most submissions will apply the obvious rule and silently corrupt a third of the corpus
without ever noticing. The differentiator is not the fix — it is having measured the
problem before writing the parser.

HOW DID YOU VERIFY IT WITHOUT GROUND TRUTH?
An invariant rather than examples. Every raw entry index must be consumed exactly once
across all reconstructed claims. 10,578 in, 10,578 out, zero gaps. That property holds
regardless of what the data looks like next week, which example-based tests cannot
promise.

ISN'T THE PARSER OVERFITTED TO THIS DATASET?
It was, initially — the cancellation-marker regex enumerated the five shapes found in
this corpus. That was rewritten to be separator-agnostic. Tests now cover formats that
appear NOWHERE in this sample: "1)" numbering, em-dash ranges, comma separators,
uppercase CANCELED, 120-claim patents, all-fragment input. Generalising it also earned
its keep — 14 entries used "N)" format and became correctly-numbered instead of inferred.

WHAT ARE THE REMAINING LIMITS?
Some source claims are truncated upstream ("wherein SA", "a land ratio La") — subscripts
and formulae were lost in the original XML parse. Merged text there reads oddly, but that
is inherited damage, not a reconstruction error. It is flagged via status, not repaired,
because inventing text in a legal document is worse than showing it damaged.

WHY CAN'T YOU PARSE THE FULL CLASSIFICATION HIERARCHY?
Codes are packed: B60B1110FI. Section (B), class (B60) and subclass (B60B) parse cleanly
from the first four characters. Group/subgroup does NOT — B60B1110FI could be 11/10 or
1/110, and there is no delimiter. So we deliberately stop at subclass and use prefix
matching, which is unambiguous and is all the task requires. Guessing would produce
filters that silently return the wrong patents.
""")


def slide_results(prs):
    s, y = new_slide(prs, "04", "Results",
                     "The benchmark lied twice before it told the truth")

    table(s, M, y, CONTENT_W,
          ["#", "Test collection", "Overlap", "Conclusion"],
          [["1", "Patent's own abstract as query", "55.4%", "BM25 wins"],
           ["2", "Queries paraphrased by an LLM", "36.9%", "Hybrid ahead, not significant"],
           ["3", "TREC pooling + graded assessor", "—", "Hybrid wins, p = 0.0003"]],
          col_w=[0.4, 4.2, 1.2, 4.2], highlight=2, row_h=Inches(0.3))
    y += Inches(1.32)

    tf = textbox(s, M, y, CONTENT_W, Inches(0.72))
    rich(tf, [("Attempt 1 ", True, INK, SANS),
              ("used each patent's own abstract — same drafter, same vocabulary, so 55% "
               "of query words appear verbatim in the target and the benchmark rewarded "
               "near-duplicate detection. ", False, INK2, SANS),
              ("Attempt 3 ", True, INK, SANS),
              ("found the real flaw: qrels marked only the query's own patent relevant, "
               "so when hybrid correctly returned five different flexible-spoke patents, "
               "four scored as errors. Pooling surfaced 1,272 cross-patent relevant "
               "records the earlier sets counted as mistakes.", False, INK2, SANS)],
         size=12, first=True, space_after=0)
    y += Inches(0.86)

    table(s, M, y, CONTENT_W,
          ["System", "recall@10", "recall@50", "nDCG@10", "P@5", "P50 latency"],
          [["BM25", "0.205", "0.405", "0.527", "0.858", "4.4 ms"],
           ["Dense", "0.253", "0.463", "0.584", "0.930", "310 ms"],
           ["Hybrid", "0.223", "0.616", "0.590", "0.885", "322 ms"]],
          col_w=[2.2, 1.4, 1.4, 1.4, 1.2, 1.6], highlight=2, row_h=Inches(0.3))
    y += Inches(1.32)

    cw = (CONTENT_W - Inches(0.3)) / 2
    mono_block(s, M, y, cw, Inches(0.9), [
        ("paired bootstrap, n=80", MUTED),
        ("recall@50  +0.2109  p=0.0001 *   +52%", GOOD),
        ("ndcg@10    +0.0627  p=0.0003 *", GOOD),
    ])
    mono_block(s, M + cw + Inches(0.3), y, cw, Inches(0.9), [
        ("+ trained reranker, 22 held-out queries", MUTED),
        ("ndcg@10    0.633 -> 0.753   p=0.0001 *", GOOD),
        ("recall@50  0.616 -> 0.616   p=1.0000  <- correct", SIGNAL),
    ])
    y += Inches(1.04)

    callout(s, M, y, CONTENT_W, Inches(0.62), None,
            "That last line is a correctness proof, not a null result: a reranker "
            "reorders the candidates it was given and cannot change which documents "
            "retrieval found. Separately — of dense search's 320 ms, 315 ms is the "
            "hosted embedding API and only 5.9 ms is OpenSearch.", kind="good")

    notes(s, """
WHAT DO THE METRICS MEAN?
- success@k — fraction of queries with at least one relevant hit in the top k. "Did the
  search find anything useful?" Not capped by how many relevant docs exist, so it reads
  on its own — but it saturates at 1.000 on a narrow corpus, which happened here.
- recall@k — share of ALL relevant documents retrieved. MUST be read against its ceiling:
  with 24 relevant documents, recall@10 cannot exceed 10/24 = 0.42. It is not scored
  against 1.0. Every report records recall_ceiling@k beside it.
- recall@50 is the metric for a FIRST-STAGE retriever: a reranker can reorder what was
  found but can never recover a claim that was never retrieved.
- precision@k — share of the top k that is relevant. The attorney's-time metric.
- MRR@k — reciprocal rank of the first relevant hit. How fast you see something useful.
- nDCG@k — the only one using the graded 0-3 scale, discounted by position. Perfect
  ordering = 1.0. THE headline for a reranker, because reranking changes order without
  changing membership.
- p-value — paired bootstrap, 10,000 resamples. Both systems run the same queries so
  scores are paired. Below 0.05 conventionally means unlikely to be chance.
  Distribution-free, because IR metrics are bounded and skewed, not normal.

ARE THESE NUMBERS GOOD? WHAT WOULD YOU NOT CLAIM?
The COMPARISON is trustworthy: hybrid beats BM25 on recall@50 by 52% at p=0.0001, on
identical qrels, so pooling bias affects both equally.
The ABSOLUTE values are flattering, and I would say so before a reviewer finds it:
 (1) the pool was built from these systems' own output, so almost every judged-relevant
     document was retrieved by something — standard TREC pooling bias;
 (2) the assessor is an LLM and graded generously (1390 grade-1, 1272 grade-2, 604
     grade-3);
 (3) in a 640-patent corpus that is entirely wheels and tyres, almost anything is
     somewhat related — which is why success@10 hit 1.000 for all three and stopped
     discriminating.
I would quote nDCG@10 and the p-values. I would NOT quote P@5 = 1.000 as evidence.

WHY LEARNING-TO-RANK RATHER THAN FINE-TUNING A CROSS-ENCODER?
Trains in seconds on CPU, scores in microseconds instead of ~90 ms, needs no model
download, and is what production search actually does — OpenSearch ships an LTR plugin
for exactly this. The trade: it only sees the features (ranks, scores, retriever
agreement, record type, overlap), never the candidate text.
Split is BY PATENT, not by row — records from one patent share vocabulary, so a row-level
split leaks the answer. Top features: vector reciprocal rank 0.285, fused score 0.243,
BM25 reciprocal rank 0.136 — the model largely learned how much to trust the semantic
ranking relative to the lexical one.

WHAT IS TREC POOLING?
Run every system, take the top-15 each, union and dedupe (~26 candidates per query),
judge each 0-3 against an examiner rubric, treat anything unjudged as non-relevant. The
standard method since 1992. 2,095 judgements, 235 API calls, cached on disk so re-running
is free and human labels can replace the LLM's wholesale.
""")


def slide_scale(prs):
    s, y = new_slide(prs, "05", "Part 2 · Implementation at scale",
                     "From 640 patents to 10 million")

    # ---- pipeline diagram
    label_box(s, M, y + Inches(0.28), Inches(1.5), Inches(0.6), "USPTO", "weekly bulk")
    arrow(s, M + Inches(1.58), y + Inches(0.52), Inches(0.3))
    label_box(s, M + Inches(1.96), y + Inches(0.28), Inches(1.4), Inches(0.6), "S3",
              "immutable", fill=PAPER)
    arrow(s, M + Inches(3.44), y + Inches(0.52), Inches(0.3))
    label_box(s, M + Inches(3.82), y + Inches(0.28), Inches(1.5), Inches(0.6), "Queue",
              "1 msg/patent", accent=DRAFT)

    wx = M + Inches(5.62)
    workers = ["Parser", "Reconstruct", "Embed (GPU)", "Index"]
    for i, w in enumerate(workers):
        label_box(s, wx + i * Inches(1.42), y, Inches(1.3), Inches(0.44), w, None,
                  title_size=10)
    arrow(s, wx - Inches(0.26), y + Inches(0.2), Inches(0.22))

    # Bus joining the four workers, with drops to the two stores they write to.
    bus_y = y + Inches(0.52)
    bus = box(s, wx + Inches(0.6), bus_y, Inches(4.6), Pt(1.25), fill=DRAFT, line=None)
    bus.shadow.inherit = False
    for i in range(4):
        tick = box(s, wx + Inches(0.65) + i * Inches(1.42), y + Inches(0.44),
                   Pt(1.25), Inches(0.08), fill=DRAFT, line=None)
        tick.shadow.inherit = False
    down_arrow(s, wx + Inches(1.55), bus_y, Inches(0.2), color=SIGNAL)
    down_arrow(s, wx + Inches(4.1), bus_y, Inches(0.2), color=DRAFT)

    label_box(s, wx + Inches(0.62), y + Inches(0.76), Inches(2.0), Inches(0.44),
              "Postgres", "ingestion_jobs", fill=SIGNAL_SOFT, line=SIGNAL, title_size=10)
    label_box(s, wx + Inches(3.05), y + Inches(0.76), Inches(2.15), Inches(0.44),
              "OpenSearch", "sharded", fill=DRAFT_SOFT, line=DRAFT, title_size=10)

    tf = textbox(s, wx + Inches(0.62), y + Inches(1.24), Inches(4.6), Inches(0.2))
    para(tf, "status writes                          indexed records",
         size=8, color=MUTED, font=MONO, space_after=0, first=True)
    y += Inches(1.56)

    sw = (CONTENT_W - Inches(0.6)) / 3
    stat(s, M, y, sw, "293M", "searchable records at 10M patents (29.3 per patent, measured)")
    stat(s, M + sw + Inches(0.3), y, sw, "~$1.4k",
         "one-time embedding backfill · ~$350 self-hosted")
    stat(s, M + 2 * (sw + Inches(0.3)), y, sw, "~$2.6k",
         "per month steady state · ~$700 optimised")
    y += Inches(1.0)

    cw = (CONTENT_W - Inches(0.3)) / 2
    tf = textbox(s, M, y, cw, Inches(1.3))
    para(tf, "Vector storage dominates — attack it first", size=12.5, bold=True,
         color=INK, font=SANS, space_after=5, first=True)
    for t in ["293M × 1536 dims × 4 bytes = 1.8 TB, and HNSW wants it resident",
              "Truncate to 512 dims (Matryoshka) -> 600 GB",
              "Quantise to int8 -> 150 GB",
              "Skip description passages — 38% of records, least queried"]:
        para(tf, "— " + t, size=11, color=INK2, font=SANS, space_after=3, line=1.15)

    tf = textbox(s, M + cw + Inches(0.3), y, cw, Inches(1.3))
    para(tf, "Proof of concept, on all 640 patents", size=12.5, bold=True,
         color=INK, font=SANS, space_after=5, first=True)
    for t in ["Idempotent — re-running enqueues 0, skips 640",
              "Versioned — bumping parser version re-queues all 640",
              "Crash-safe — expired leases return to PENDING",
              "Retry-tracked — failures record the error and count"]:
        para(tf, "— " + t, size=11, color=INK2, font=SANS, space_after=3, line=1.15)
    y += Inches(1.42)

    callout(s, M, y, CONTENT_W, Inches(0.6), None,
            "The filter finding inverts at scale: filters cost BM25 +69% here because at "
            "18.7K records there is nothing to save. At 10M patents B60B selects ~0.3% of "
            "the corpus, and filtering becomes the dominant optimisation.", kind="draft")

    notes(s, """
HOW DID YOU GET THESE COST NUMBERS?
Extrapolated from measurements on the 640-patent sample, not guessed. Measured basis:
29.3 records per patent, 24K characters per patent, ~6,900 embedding tokens per patent,
2.5s record build, 2.2s bulk index. Scaled to 10^7 patents: 293M records, ~240 GB raw
text, ~69B embedding tokens. At $0.02/1M tokens that is ~$1,400 for a full hosted
backfill; self-hosting MiniLM on 8x g5.xlarge for ~40h is ~$350, so self-host the
backfill and use the hosted API for the weekly delta where operational simplicity wins.

WHAT BREAKS FIRST AT SCALE?
Vector index memory. 293M x 1536 dims x 4 bytes = 1.8 TB of float32, and HNSW wants it
resident. That is the binding constraint, not CPU. Three levers in order of value:
dimension reduction (Matryoshka to 512d, already wired as --dimensions), int8
quantisation (4x), and not embedding description passages at all (38% of records, least
queried). All three together bring the cluster from ~$2,600/mo to ~$700/mo.

WHAT ELSE IS HARD AT SCALE?
- Reindexing on model change: swapping embedding models invalidates 293M vectors. Version
  the index, build the new one alongside, alias-swap. Never reindex in place.
- Claim reconstruction drift: upstream XML format changes silently corrupt claims. The
  invariant test runs per release, and the number_inferred rate is an alarm.
- Filter selectivity: B60B matches 47% here but a narrow CPC code may match 0.001% of
  10M. Pre-filter in traversal, fall back to lexical-only below a cardinality threshold.
- Hot shards: recent filings are queried far more. Time-based index tiering.

WHY SQLITE FOR THE PROOF OF CONCEPT?
It is stdlib, needs no daemon or container, and has the same semantics we want from
Postgres — transactional claim via BEGIN IMMEDIATE, unique primary key, atomic status
transitions. Swapping the DSN is the only change needed for production. The POC
demonstrates idempotency, versioned reprocessing, crash recovery via lease expiry, and
retry accounting, all on the real 640-patent corpus.

WHAT IS STILL UNPROVEN?
640 patents is a POC. Behaviour at 10M is argued, not measured. build_index.py holds all
embeddings in memory before writing — fine here, won't scale; the queue-driven worker
path is the one that does. The API has no auth or rate limiting, correct for localhost
and wrong for anything exposed.
""")


def build(out: Path) -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    slide_problem(prs)
    slide_architecture(prs)
    slide_reconstruction(prs)
    slide_results(prs)
    slide_scale(prs)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    n = len(prs.slides._sldIdLst)
    noted = sum(1 for s in prs.slides if s.has_notes_slide)
    print(f"wrote {out}")
    print(f"  {n} slides, {noted} with speaker notes")


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("-o", "--out", type=Path,
                    default=Path(__file__).resolve().parents[1] / "docs" / "thinkstruct-deck.pptx")
    build(ap.parse_args().out)


